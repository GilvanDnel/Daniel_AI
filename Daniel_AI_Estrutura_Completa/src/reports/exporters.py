"""Report export helpers for Daniel AI."""

from __future__ import annotations

from io import BytesIO
from typing import Any

import pandas as pd
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


def export_dataframe_csv(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8-sig")


def export_dataframe_excel(df: pd.DataFrame) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="resultado")
    return output.getvalue()


def export_summary_pdf(title: str, summary: dict[str, Any]) -> bytes:
    output = BytesIO()
    document = SimpleDocTemplate(output, pagesize=A4, title=title)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]

    for key, value in summary.items():
        label = key.replace("_", " ").title()
        story.append(Paragraph(f"<b>{label}:</b> {value}", styles["BodyText"]))
        story.append(Spacer(1, 6))

    document.build(story)
    return output.getvalue()


def export_summary_pptx(title: str, summary: dict[str, Any]) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Resumo gerado pelo Daniel AI"

    content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    content_slide.shapes.title.text = "Principais Indicadores"
    body = content_slide.placeholders[1].text_frame
    body.clear()

    for key, value in summary.items():
        paragraph = body.add_paragraph()
        paragraph.text = f"{key.replace('_', ' ').title()}: {value}"
        paragraph.level = 0

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()
