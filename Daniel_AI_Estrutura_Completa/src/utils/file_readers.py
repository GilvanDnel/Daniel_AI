"""Text extractors for supported document formats."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def read_pdf(file_path: str | Path) -> str:
    reader = PdfReader(str(file_path))
    parts = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"[Página {index}]\n{text}")
    return "\n\n".join(parts)


def read_docx(file_path: str | Path) -> str:
    doc = Document(str(file_path))
    parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_pptx(file_path: str | Path) -> str:
    presentation = Presentation(str(file_path))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            parts.append(f"[Slide {index}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def read_markdown(file_path: str | Path) -> str:
    return Path(file_path).read_text(encoding="utf-8")


def read_document(file_path: str | Path) -> str:
    file_path = Path(file_path)
    readers = {
        ".pdf": read_pdf,
        ".docx": read_docx,
        ".pptx": read_pptx,
        ".md": read_markdown,
        ".txt": read_markdown,
    }
    reader = readers.get(file_path.suffix.lower())
    if reader is None:
        raise ValueError(f"Formato não suportado para leitura textual: {file_path.suffix}")
    return reader(file_path)
