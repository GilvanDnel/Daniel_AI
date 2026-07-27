"""
file_readers.py

Responsável por extrair texto de diferentes formatos de arquivo:
PDF, DOCX, PPTX. CSV e XLSX são tratados separadamente pelo módulo
de analytics (pandas), não por aqui.
"""

from __future__ import annotations

from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def read_pdf(file_path: str | Path) -> str:
    """Extrai todo o texto de um arquivo PDF."""
    reader = PdfReader(file_path)
    texto = []
    for i, page in enumerate(reader.pages):
        conteudo = page.extract_text() or ""
        if conteudo.strip():
            texto.append(f"[Página {i + 1}]\n{conteudo}")
    return "\n\n".join(texto)


def read_docx(file_path: str | Path) -> str:
    """Extrai todo o texto de um arquivo DOCX (parágrafos e tabelas)."""
    doc = Document(file_path)
    partes = []

    for paragrafo in doc.paragraphs:
        if paragrafo.text.strip():
            partes.append(paragrafo.text)

    for tabela in doc.tables:
        for linha in tabela.rows:
            celulas = [celula.text.strip() for celula in linha.cells]
            partes.append(" | ".join(celulas))

    return "\n".join(partes)


def read_pptx(file_path: str | Path) -> str:
    """Extrai todo o texto de um arquivo PPTX (slide a slide)."""
    prs = Presentation(file_path)
    partes = []

    for i, slide in enumerate(prs.slides):
        textos_slide = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    texto_par = "".join(run.text for run in paragrafo.runs)
                    if texto_par.strip():
                        textos_slide.append(texto_par)
        if textos_slide:
            partes.append(f"[Slide {i + 1}]\n" + "\n".join(textos_slide))

    return "\n\n".join(partes)


def read_markdown(file_path: str | Path) -> str:
    """Extrai o texto de um arquivo Markdown (.md), lido como texto puro."""
    return Path(file_path).read_text(encoding="utf-8")


def read_document(file_path: str | Path) -> str:
    """
    Detecta a extensão do arquivo e chama o leitor apropriado.
    Levanta ValueError se o formato não for suportado.
    """
    file_path = Path(file_path)
    extensao = file_path.suffix.lower()

    leitores = {
        ".pdf": read_pdf,
        ".docx": read_docx,
        ".pptx": read_pptx,
        ".md": read_markdown,
    }

    if extensao not in leitores:
        raise ValueError(
            f"Formato '{extensao}' não suportado para leitura de texto. "
            f"Formatos aceitos: {', '.join(leitores.keys())}"
        )

    return leitores[extensao](file_path)
