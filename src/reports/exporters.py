"""Report export helpers for Daniel AI."""

from __future__ import annotations

import html
from io import BytesIO
from typing import Any

import pandas as pd
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


def export_dataframe_csv(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8-sig")


def export_dataframe_excel(df: pd.DataFrame) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="resultado")
    return output.getvalue()


def _analysis_parts(analysis_or_summary: Any) -> tuple[dict[str, Any], list[str], list[str]]:
    if isinstance(analysis_or_summary, dict):
        return analysis_or_summary, [], []
    return (
        getattr(analysis_or_summary, "summary", {}),
        getattr(analysis_or_summary, "insights", []),
        getattr(analysis_or_summary, "recommendations", []),
    )


def export_summary_pdf(title: str, analysis_or_summary: Any) -> bytes:
    summary, insights, recommendations = _analysis_parts(analysis_or_summary)
    output = BytesIO()
    document = SimpleDocTemplate(output, pagesize=A4, title=title)
    styles = getSampleStyleSheet()

    safe_title = html.escape(title)
    story = [Paragraph(safe_title, styles["Title"]), Spacer(1, 12)]

    if insights:
        story.append(Paragraph("Leitura Executiva", styles["Heading2"]))
        for insight in insights:
            story.append(Paragraph(f"- {html.escape(str(insight))}", styles["BodyText"]))
        story.append(Spacer(1, 10))

    if recommendations:
        story.append(Paragraph("Recomendações", styles["Heading2"]))
        for recommendation in recommendations:
            story.append(Paragraph(f"- {html.escape(str(recommendation))}", styles["BodyText"]))
        story.append(Spacer(1, 10))

    story.append(Paragraph("KPIs", styles["Heading2"]))
    for key, value in summary.items():
        label = html.escape(key.replace("_", " ").title())
        if isinstance(value, list):
            formatted_val = html.escape(", ".join(map(str, value)))
        else:
            formatted_val = html.escape(str(value))
        story.append(Paragraph(f"<b>{label}:</b> {formatted_val}", styles["BodyText"]))
        story.append(Spacer(1, 5))

    document.build(story)
    return output.getvalue()


def export_summary_pptx(title: str, analysis_or_summary: Any) -> bytes:
    summary, insights, recommendations = _analysis_parts(analysis_or_summary)
    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Resumo executivo gerado pelo Daniel AI"

    if insights:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Leitura Executiva"
        body = slide.placeholders[1].text_frame
        body.clear()
        for insight in insights[:5]:
            paragraph = body.add_paragraph()
            paragraph.text = str(insight)
            paragraph.level = 0

    if recommendations:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Recomendações"
        body = slide.placeholders[1].text_frame
        body.clear()
        for recommendation in recommendations[:5]:
            paragraph = body.add_paragraph()
            paragraph.text = str(recommendation)
            paragraph.level = 0

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Principais Indicadores"
    body = slide.placeholders[1].text_frame
    body.clear()
    for key, value in summary.items():
        paragraph = body.add_paragraph()
        if isinstance(value, list):
            formatted_val = ", ".join(map(str, value))
        else:
            formatted_val = str(value)
        paragraph.text = f"{key.replace('_', ' ').title()}: {formatted_val}"
        paragraph.level = 0

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def export_chat_history_txt(messages: list[dict]) -> bytes:
    """Export conversation transcript as a clean UTF-8 text file."""
    lines = [
        "==========================================",
        "   HISTÓRICO DE CONVERSA - DANIEL AI",
        "   DNEL SOM Serviços Inteligentes",
        "==========================================",
        "",
    ]
    for msg in messages:
        role = "Daniel (Assistente)" if msg.get("role") == "assistant" else "Usuário"
        content = msg.get("content", "")
        lines.append(f"[{role}]")
        lines.append(content)
        lines.append("-" * 40)

    return "\n".join(lines).encode("utf-8")


def export_chat_history_pdf(messages: list[dict]) -> bytes:
    """Export conversation transcript as an executive PDF document."""
    output = BytesIO()
    document = SimpleDocTemplate(output, pagesize=A4, title="Histórico da Conversa - Daniel AI")
    styles = getSampleStyleSheet()

    story = [
        Paragraph("Histórico de Atendimento - Daniel AI", styles["Title"]),
        Spacer(1, 14),
    ]

    for msg in messages:
        role = "Daniel (Assistente)" if msg.get("role") == "assistant" else "Usuário"
        content = msg.get("content", "")
        header_text = f"<b>{html.escape(role)}</b>"
        story.append(Paragraph(header_text, styles["Heading3"]))
        clean_content = html.escape(content).replace("\n", "<br/>")
        story.append(Paragraph(clean_content, styles["BodyText"]))
        story.append(Spacer(1, 8))

    document.build(story)
    return output.getvalue()

